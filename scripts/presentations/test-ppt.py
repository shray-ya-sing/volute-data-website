import json
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_presentation(data):
    company_name = data.get('companyName', 'Unknown Company')
    overview = data.get('overview', 'No overview provided.')
    key_points = data.get('keyPoints', [])

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === TITLE SLIDE ===
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(31, 78, 120)  # Dark blue

    # Title
    title = slide.shapes.title
    title.text = company_name
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Company Overview"
        subtitle.text_frame.paragraphs[0].font.size = Pt(32)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # === CONTENT SLIDE ===
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)

    # Title
    title_shape = slide.shapes.title
    title_shape.text = f"{company_name} Overview"
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 78, 120)

    # Body
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = overview
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].space_after = Pt(20)

    # Add key points with bullets
    for point in key_points:
        p = tf.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(16)
        p.space_after = Pt(12)

    # Save
    filename = company_name.replace(" ", "_").replace("/", "-")
    output_path = f'{filename}_overview.pptx'
    prs.save(output_path)
    return os.path.abspath(output_path)

if __name__ == "__main__":
    # Check if JSON string is passed via command line, else use dummy data
    if len(sys.argv) > 1:
        raw_data = json.loads(sys.argv[1])
    else:
        # Sample data for local testing
        raw_data = {
            "companyName": "Test BioTech",
            "overview": "A leading innovator in neural engineering and AI-driven drug discovery.",
            "keyPoints": [
                "Founded in 2022",
                "Raised $50M in Series A",
                "Phase II clinical trials underway"
            ]
        }
    
    path = create_presentation(raw_data)
    print(f"Success! File created at: {path}")
    