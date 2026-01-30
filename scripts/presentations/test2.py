from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Add blank slide
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Define colors
NAVY_BLUE = RGBColor(0, 32, 91)
LIGHT_GRAY = RGBColor(242, 242, 242)
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(89, 89, 89)
MEDIUM_GRAY = RGBColor(166, 166, 166)

def add_textbox(slide, left, top, width, height, text, font_size, bold=False, 
                color=BLACK, align=PP_ALIGN.LEFT, bg_color=None, italic=False):
    """Helper function to add text box"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    
    p = text_frame.paragraphs[0]
    p.alignment = align
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    
    if bg_color:
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    
    return textbox

def add_rectangle(slide, left, top, width, height, fill_color):
    """Helper function to add rectangle"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = fill_color
    return shape

def set_cell_background(cell, fill_color):
    """Set cell background color"""
    try:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color
    except:
        pass

def set_cell_border(cell, border_color, border_width=Pt(1)):
    """Set cell border"""
    try:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Add borders
        for border_name in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
            ln = parse_xml(f'<{border_name} w="{int(border_width)}" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"><a:solidFill><a:srgbClr val="{border_color.rgb_hex[2:]}"/></a:solidFill></{border_name}>')
            tcPr.append(ln)
    except:
        pass

# Title
title_box = add_textbox(slide, Inches(0.5), Inches(0.3), Inches(9), Inches(0.6),
                         "Executive Summary", 48, bold=True, color=BLACK)

# Business Overview Section
box_width = Inches(4.5)
box_height = Inches(0.4)

# Business Overview Header
overview_header = add_rectangle(slide, Inches(0.4), Inches(1.1), box_width, box_height, NAVY_BLUE)
overview_title = add_textbox(slide, Inches(0.4), Inches(1.1), box_width, box_height,
                              "Business Overview", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
overview_title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Business Overview Content
content_top = Inches(1.6)
bullet_points = [
    "Nintendo (TYO:7974) is a best-in-class international video game developer and video game console manufacturer, headquartered in Kyoto, Japan",
    "The company is a global market leader, having contributed heavily to the industry's inception in the 1980s and experiencing multi-decade success ever since",
    "Nintendo's stock is currently trading at a discount to its intrinsic value as well as relative to peers, making now an opportune time to invest in the company",
    "High grow rates since 2016 are driven largely by the company's release of the Nintendo Switch in 2017"
]

for i, bullet in enumerate(bullet_points):
    textbox = slide.shapes.add_textbox(Inches(0.5), content_top + Inches(i * 0.5), 
                                       Inches(4.3), Inches(0.45))
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = bullet
    p.font.size = Pt(9)
    p.font.color.rgb = BLACK
    p.level = 0
    p.space_before = Pt(2)
    p.space_after = Pt(2)

# Historical Financials Section
fin_left = Inches(5.1)
fin_header = add_rectangle(slide, fin_left, Inches(1.1), box_width, box_height, NAVY_BLUE)
fin_title = add_textbox(slide, fin_left, Inches(1.1), box_width, box_height,
                        "Historical Financials", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
fin_title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Financial Table using actual table object
table_left = fin_left + Inches(0.05)
table_top = Inches(1.6)
table_width = box_width - Inches(0.1)
table_height = Inches(2.0)

# Create table (11 rows x 7 columns)
rows = 11
cols = 7
table_shape = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height)
table = table_shape.table

# Set column widths
col_widths = [Inches(1.0), Inches(0.58), Inches(0.58), Inches(0.58), Inches(0.58), Inches(0.58), Inches(0.48)]
for col_idx in range(cols):
    table.columns[col_idx].width = col_widths[col_idx]

# Table data
table_data = [
    ["", "2015A", "2016A", "2017A", "2018A", "LTM", "CAGR"],
    ["Revenue", "$4,430.3", "$3,339.8", "$9,188.8", "$10,900.7", "$11,175.2", "26.0%"],
    ["YoY Growth", "", "(24.6%)", "175.1%", "18.6%", "2.5%", ""],
    ["Gross Profit", "$1,871.4", "$1,622.3", "$3,419.3", "$4,432.5", "$4,654.6", "25.6%"],
    ["Margin", "", "45.6%", "37.2%", "40.7%", "41.7%", ""],
    ["EBITDA", "$365.9", "$222.2", "$1,485.3", "$2,286.1", "$2,377.7", "59.7%"],
    ["YoY Growth", "", "(39.3%)", "568.3%", "53.8%", "4.1%", ""],
    ["Margin", "", "6.7%", "16.2%", "21.0%", "21.3%", ""],
    ["Note: R&D", "$526.9", "$591.6", "$525.6", "$583.6", "$583.6", ""],
    ["EBITDA Excl. R&D", "$892.8", "$813.8", "$2,010.9", "$2,868.7", "$2,961.4", "35.0%"],
    ["Margin", "", "24.4%", "21.9%", "26.3%", "26.5%", ""]
]

# Populate table
for row_idx in range(rows):
    for col_idx in range(cols):
        cell = table.cell(row_idx, col_idx)
        cell.text = table_data[row_idx][col_idx]
        
        # Format cell
        text_frame = cell.text_frame
        text_frame.margin_left = Pt(3)
        text_frame.margin_right = Pt(3)
        text_frame.margin_top = Pt(2)
        text_frame.margin_bottom = Pt(2)
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Format paragraph
        if text_frame.paragraphs:
            p = text_frame.paragraphs[0]
            
            # Set alignment
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.RIGHT
            
            # Set font properties
            is_header = row_idx == 0
            is_label = col_idx == 0
            is_italic_row = row_idx in [2, 4, 7, 10]
            
            p.font.size = Pt(7 if not is_header else 8)
            p.font.bold = is_header or is_label
            p.font.color.rgb = BLACK
            
            if is_italic_row and col_idx == 0:
                p.font.italic = True
        
        # Set cell background for header row
        if row_idx == 0:
            set_cell_background(cell, LIGHT_GRAY)

# Deal History Section
deal_left = Inches(0.4)
deal_top = Inches(3.8)
deal_header = add_rectangle(slide, deal_left, deal_top, box_width, box_height, NAVY_BLUE)
deal_title = add_textbox(slide, deal_left, deal_top, box_width, box_height,
                          "Deal History", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
deal_title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Deal History Content (placeholder text)
deal_content = slide.shapes.add_textbox(
    deal_left + Inches(0.1), 
    deal_top + Inches(0.5), 
    box_width - Inches(0.2), 
    Inches(2.4)
)
deal_text_frame = deal_content.text_frame
deal_text_frame.word_wrap = True
deal_text_frame.margin_left = Inches(0.1)
deal_text_frame.margin_right = Inches(0.1)
deal_text_frame.margin_top = Inches(0.1)
deal_text_frame.margin_bottom = Inches(0.1)

# Add placeholder paragraph
p = deal_text_frame.paragraphs[0]
p.text = "[DEAL HISTORY PLACEHOLDER]\n\n"
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY
p.font.italic = True

# Add additional paragraphs for formatting
for i in range(3):
    p = deal_text_frame.add_paragraph()
    p.text = "• [Deal details to be populated here]"
    p.font.size = Pt(9)
    p.font.color.rgb = DARK_GRAY
    p.font.italic = True
    p.space_after = Pt(6)

# Recent News Section
news_left = fin_left
news_header = add_rectangle(slide, news_left, deal_top, box_width, box_height, NAVY_BLUE)
news_title = add_textbox(slide, news_left, deal_top, box_width, box_height,
                           "Recent News", 16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
news_title.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Recent News Content (placeholder text)
news_content = slide.shapes.add_textbox(
    news_left + Inches(0.1), 
    deal_top + Inches(0.5), 
    box_width - Inches(0.2), 
    Inches(2.4)
)
news_text_frame = news_content.text_frame
news_text_frame.word_wrap = True
news_text_frame.margin_left = Inches(0.1)
news_text_frame.margin_right = Inches(0.1)
news_text_frame.margin_top = Inches(0.1)
news_text_frame.margin_bottom = Inches(0.1)

# Add placeholder paragraph
p = news_text_frame.paragraphs[0]
p.text = "[RECENT NEWS PLACEHOLDER]\n\n"
p.font.size = Pt(10)
p.font.color.rgb = DARK_GRAY
p.font.italic = True

# Add additional paragraphs for formatting
for i in range(3):
    p = news_text_frame.add_paragraph()
    p.text = "• [News item to be populated here]"
    p.font.size = Pt(9)
    p.font.color.rgb = DARK_GRAY
    p.font.italic = True
    p.space_after = Pt(6)

# Page number
add_textbox(slide, Inches(9.2), Inches(7.0), Inches(0.3), Inches(0.3),
            "4", 14, color=BLACK, align=PP_ALIGN.RIGHT)

# Save presentation
prs.save('nintendo_executive_summary_updated.pptx')
print("Presentation created successfully!")